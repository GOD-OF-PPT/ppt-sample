import io
import re
import numpy as np
import fitz  # PyMuPDF
from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

app = Flask(__name__, static_folder=".", static_url_path="")

# ── Question detection ────────────────────────────────────────────────────────
Q_PATTERN = re.compile(
    r"^[\s　]*(?:第\s*(\d+)\s*[题题]|(\d+)\s*[\.、．\.\uff0e])",
    re.MULTILINE,
)
RENDER_SCALE        = 2.0   # hi-res output (144 dpi)
OCR_SCALE           = 1.5   # lower scale for OCR question-detection pass
TOP_MARGIN_PTS      = 6     # pts above detected question start to include
OCR_Y_THRESHOLD_PTS = 200   # max PDF pts from page-top for question-num search
MIN_SEGMENT_H_PTS   = 150   # skip page segments shorter than this (whitespace strips)

# ── Slide layout ─────────────────────────────────────────────────────────────
SLIDE_W_IN  = 13.33   # widescreen landscape
SLIDE_H_IN  = 7.5
TITLE_H_IN  = 0.55
MARGIN_W_IN = 0.0
MARGIN_H_IN = 0.0


# ─────────────────────────────────────────────────────────────────────────────
# QUESTION DETECTION
# ─────────────────────────────────────────────────────────────────────────────

def _block_q_num(block, page_width):
    if block[6] != 0:
        return None
    if block[0] > page_width * 0.40:
        return None
    text = block[4].strip()
    m = Q_PATTERN.match(text)
    if not m:
        return None
    return int(m.group(1) or m.group(2))


def detect_questions(doc):
    """Text-layer based question detection."""
    questions, seen = [], set()
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        pw = page.rect.width
        for block in sorted(page.get_text("blocks"), key=lambda b: b[1]):
            qnum = _block_q_num(block, pw)
            if qnum and qnum not in seen:
                seen.add(qnum)
                questions.append({
                    "num": qnum,
                    "page_idx": page_idx,
                    "y_start": max(0.0, block[1] - TOP_MARGIN_PTS),
                })
    questions.sort(key=lambda q: (q["page_idx"], q["y_start"]))
    return questions


def detect_questions_ocr(doc):
    """OCR-based fallback for image-only PDFs."""
    try:
        import pytesseract
    except ImportError:
        return []

    questions, seen = [], set()
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        pix = page.get_pixmap(matrix=fitz.Matrix(OCR_SCALE, OCR_SCALE), alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        data = pytesseract.image_to_data(
            img, output_type=pytesseract.Output.DICT, lang="eng"
        )
        n = len(data["text"])
        for i in range(n):
            raw  = data["text"][i].strip()
            word = re.sub(r"[\.、．,;:]+$", "", raw)
            if not re.match(r"^\d{1,2}$", word):
                continue
            qnum = int(word)
            if qnum < 1 or qnum > 99 or qnum in seen:
                continue
            x_px, y_px = data["left"][i], data["top"][i]
            if x_px > pix.width * 0.20:
                continue
            y_pdf = y_px / OCR_SCALE
            if y_pdf > OCR_Y_THRESHOLD_PTS:
                continue
            if int(data["conf"][i]) < 35 or data["word_num"][i] != 1:
                continue
            blk, ln = data["block_num"][i], data["line_num"][i]
            nexts = [
                data["text"][j].strip()
                for j in range(i + 1, min(i + 6, n))
                if data["block_num"][j] == blk and data["line_num"][j] == ln
                   and data["text"][j].strip()
            ]
            if not nexts or not re.match(r"[A-Za-z\u4e00-\u9fff]", nexts[0]):
                continue
            seen.add(qnum)
            questions.append({
                "num": qnum,
                "page_idx": page_idx,
                "y_start": max(0.0, y_pdf - TOP_MARGIN_PTS),
            })

    questions.sort(key=lambda q: (q["page_idx"], q["y_start"]))
    return questions


# ─────────────────────────────────────────────────────────────────────────────
# RENDERING
# ─────────────────────────────────────────────────────────────────────────────

def crop_lr_whitespace(img, padding=20, threshold=245):
    """
    Crop left/right whitespace from a PIL Image independently.
    Scans each column; columns where all pixels are >= threshold are white.
    Keeps `padding` pixels of margin on each side of the detected content.
    """
    arr = np.array(img.convert("L"))   # grayscale → 2-D array (rows × cols)
    col_min = arr.min(axis=0)          # darkest pixel in each column
    content_cols = np.where(col_min < threshold)[0]
    if content_cols.size == 0:
        return img
    left  = max(0,          int(content_cols[0])  - padding)
    right = min(img.width,  int(content_cols[-1]) + padding + 1)
    return img.crop((left, 0, right, img.height))


def render_question_pages(doc, q_start, q_end):
    """
    Render each PDF page segment of a question as a separate hi-res PIL Image.
    Segments shorter than MIN_SEGMENT_H_PTS (whitespace margins) are skipped.
    Returns a non-empty list of images.
    """
    mat = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)
    sp, sy = q_start["page_idx"], q_start["y_start"]
    ep = q_end["page_idx"] if q_end else len(doc) - 1
    ey = q_end["y_start"] if q_end else None
    pages = []
    for pi in range(sp, ep + 1):
        page = doc[pi]
        pw, ph = page.rect.width, page.rect.height
        if   pi == sp and pi == ep:
            clip = fitz.Rect(0, sy, pw, ey if ey is not None else ph)
        elif pi == sp:
            clip = fitz.Rect(0, sy, pw, ph)
        elif pi == ep:
            clip = fitz.Rect(0, 0, pw, ey if ey is not None else ph)
        else:
            clip = fitz.Rect(0, 0, pw, ph)
        if clip.height < MIN_SEGMENT_H_PTS:
            continue
        pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        pages.append(crop_lr_whitespace(img))
    return pages or [Image.new("RGB", (100, 100), "white")]


# ─────────────────────────────────────────────────────────────────────────────
# PPTX BUILDER  — A4 portrait, screenshot tiles
# ─────────────────────────────────────────────────────────────────────────────

def _add_title_bar(slide, slide_w, title_h, label):
    tb = slide.shapes.add_textbox(0, 0, slide_w, title_h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor(0x26, 0x5F, 0xAB)
    tf = tb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = f"  {label}"
    run.font.size      = Pt(16)
    run.font.bold      = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def build_pptx(questions):
    """
    questions: list of {"num": int, "pages": [PIL.Image, ...]}
    One slide per question; page screenshots arranged side-by-side, no gaps.
    Returns BytesIO of the PPTX.
    """
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    slide_w   = prs.slide_width
    slide_h   = prs.slide_height
    title_h   = Inches(TITLE_H_IN)
    content_h = slide_h - title_h

    for item in questions:
        slide = prs.slides.add_slide(blank)
        _add_title_bar(slide, slide_w, title_h, f"第 {item['num']} 题")

        pages = item["pages"]
        n     = len(pages)
        col_w = int(slide_w / n)   # integer EMU, equal column width

        for i, img in enumerate(pages):
            img_w, img_h = img.size
            aspect = img_w / img_h

            # Scale image to fill column width; cap at content height
            disp_w = Emu(col_w)
            disp_h = Emu(int(col_w / aspect))
            if disp_h > content_h:
                disp_h = content_h
                disp_w = Emu(int(content_h * aspect))

            left = Emu(i * col_w)
            buf  = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, left, title_h, disp_w, disp_h)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ─────────────────────────────────────────────────────────────────────────────
# FLASK ROUTES
# ─────────────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return app.send_static_file("index.html")


@app.route("/convert", methods=["POST"])
def convert():
    if "pdf" not in request.files:
        return jsonify({"error": "未收到 PDF 文件"}), 400
    pdf_bytes = request.files["pdf"].read()
    if not pdf_bytes:
        return jsonify({"error": "PDF 文件为空"}), 400
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as e:
        return jsonify({"error": f"无法打开 PDF：{e}"}), 400

    questions = detect_questions(doc)
    if not questions:
        questions = detect_questions_ocr(doc)

    if questions:
        items = []
        for i, q in enumerate(questions):
            q_end  = questions[i + 1] if i + 1 < len(questions) else None
            pages  = render_question_pages(doc, q, q_end)
            items.append({"num": q["num"], "pages": pages})
    else:
        # Fallback: one slide per page
        mat = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)
        items = []
        for pi in range(len(doc)):
            page = doc[pi]
            pix  = page.get_pixmap(matrix=mat, alpha=False)
            img  = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            items.append({"num": pi + 1, "pages": [img]})

    pptx_buf = build_pptx(items)
    doc.close()

    return send_file(
        pptx_buf,
        mimetype=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
        as_attachment=True,
        download_name="exam.pptx",
    )


if __name__ == "__main__":
    app.run(debug=True, port=5000)
